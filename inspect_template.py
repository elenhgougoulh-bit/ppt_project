from pptx import Presentation


def analyze_template(template_path):
    print(f"Αναλύω το template: {template_path}...\n")

    try:
        prs = Presentation(template_path)
    except Exception as e:
        print(f"Σφάλμα κατά το άνοιγμα του αρχείου. Σιγουρέψου ότι υπάρχει: {e}")
        return

    # Διατρέχουμε όλα τα διαθέσιμα layouts (σχεδιαγράμματα)
    for i, layout in enumerate(prs.slide_layouts):
        print(f"--- Layout {i}: {layout.name} ---")

        # Διατρέχουμε τα placeholders (τα κουτιά) του κάθε layout
        for shape in layout.placeholders:
            print(
                f"  -> Placeholder ID: {shape.placeholder_format.idx} | Όνομα: {shape.name}"
            )
        print("-" * 40)


if __name__ == "__main__":
    # Βάλε εδώ το όνομα του αρχείου σου
    analyze_template("template.pptx")
