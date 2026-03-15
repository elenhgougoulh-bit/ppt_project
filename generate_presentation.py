from pptx import Presentation

# Αυτή είναι η δομή των δεδομένων που θα μας δίνει το AI αργότερα!
# Προς το παρόν τη γράφουμε εμείς με το χέρι για να τεστάρουμε το design.
presentation_data = [
    {
        "layout": 0,  # Title Slide
        "title": "Η Επανάσταση της Τεχνητής Νοημοσύνης",
        "subtitle": "Πώς τα LLMs αλλάζουν τον κόσμο",
    },
    {
        "layout": 2,  # Title and Content
        "title": "Τι είναι τα Μεγάλα Γλωσσικά Μοντέλα;",
        "content": "• Εκπαιδεύονται σε τεράστιο όγκο δεδομένων.\n• Κατανοούν και παράγουν φυσική γλώσσα.\n• Μπορούν να γράψουν κώδικα και ποιήματα!",
    },
    {
        "layout": 7,  # Section Header
        "title": "Οφέλη για τις Επιχειρήσεις",
        "content": "Μείωση κόστους και αύξηση παραγωγικότητας.",
    },
    {
        "layout": 40,  # Quote
        "title": "Steve Jobs",
        "content": '"Η καινοτομία ξεχωρίζει τον ηγέτη από τον ακόλουθο."',
    },
]


def create_presentation_from_template(template_path, output_path, data):
    print(f"Ανοίγω το template: {template_path}...")
    prs = Presentation(template_path)

    # Διαγράφουμε τυχόν υπάρχουσες διαφάνειες που ξέμειναν στο template
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    # Δημιουργούμε τις νέες διαφάνειες με βάση τα δεδομένα μας
    for slide_data in data:
        layout_num = slide_data["layout"]
        print(f"Προσθήκη διαφάνειας με Layout {layout_num}...")

        slide_layout = prs.slide_layouts[layout_num]
        slide = prs.slides.add_slide(slide_layout)

        # Γεμίζουμε τα placeholders ανάλογα με το layout
        if layout_num == 0:  # Title Slide (ID 0: Title, ID 1: Subtitle)
            slide.placeholders[0].text = slide_data["title"]
            slide.placeholders[1].text = slide_data["subtitle"]

        elif layout_num == 2:  # Title and Content (ID 0: Title, ID 1: Content)
            slide.placeholders[0].text = slide_data["title"]
            slide.placeholders[1].text = slide_data["content"]

        elif layout_num == 7:  # Section Header (ID 0: Title, ID 1: Text)
            slide.placeholders[0].text = slide_data["title"]
            slide.placeholders[1].text = slide_data["content"]

        elif layout_num == 40:  # Quote (ID 0: Title, ID 1: Content)
            slide.placeholders[0].text = slide_data["title"]
            slide.placeholders[1].text = slide_data["content"]

    prs.save(output_path)
    print(f"\nΕπιτυχία! Η παρουσίαση αποθηκεύτηκε στο: {output_path}")


if __name__ == "__main__":
    template_file = "template.pptx"
    output_file = "final_presentation.pptx"

    create_presentation_from_template(template_file, output_file, presentation_data)
