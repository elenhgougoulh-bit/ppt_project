import os
import requests
from pptx import Presentation
from ai_brain import generate_presentation_json
from dotenv import load_dotenv

load_dotenv()

UNSPLASH_KEY = os.getenv("UNSPLASH_ACCESS_KEY")


def get_unsplash_image(query):
    """Ψάχνει μια εικόνα στο Unsplash και επιστρέφει το URL της."""
    print(f"Αναζήτηση εικόνας για: '{query}'...")
    url = f"https://api.unsplash.com/search/photos?query={query}&client_id={UNSPLASH_KEY}&per_page=1"
    try:
        response = requests.get(url)
        data = response.json()
        if data["results"]:
            return data["results"][0]["urls"]["regular"]
    except Exception as e:
        print(f"Σφάλμα στο Unsplash API: {e}")
    return None


def download_image(url, filename="temp_image.jpg"):
    """Κατεβάζει την εικόνα τοπικά για να την εισάγει στο PPTX."""
    response = requests.get(url)
    if response.status_code == 200:
        with open(filename, "wb") as f:
            f.write(response.content)
        return filename
    return None


def build_presentation(topic):
    # 1. Παίρνουμε το περιεχόμενο από το AI
    slides_data = generate_presentation_json(topic)
    if not slides_data:
        print("Αποτυχία παραγωγής περιεχομένου.")
        return

    # 2. Φορτώνουμε το template
    prs = Presentation("template.pptx")

    # Διαγραφή παλιών διαφανειών αν υπάρχουν
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    # 3. Δημιουργούμε τις νέες διαφάνειες
    for data in slides_data:
        layout_num = data.get("layout", 2)  # Default στο Title and Content
        slide_layout = prs.slide_layouts[layout_num]
        slide = prs.slides.add_slide(slide_layout)

        print(f"Δημιουργία διαφάνειας: {data.get('title', 'Χωρίς Τίτλο')}")

        # Γεμίζουμε τα κλασικά Text Placeholders
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 0:  # Τίτλος
                placeholder.text = data.get("title", "")
            elif placeholder.placeholder_format.idx == 1:  # Περιεχόμενο/Υπότιτλος
                raw_content = data.get("content", data.get("subtitle", ""))

                # ΕΔΩ ΕΙΝΑΙ Η ΔΙΟΡΘΩΣΗ:
                # Αν το περιεχόμενο είναι λίστα, το ενώνουμε σε ένα κείμενο με αλλαγές γραμμής
                if isinstance(raw_content, list):
                    placeholder.text = "\n".join(raw_content)
                else:
                    placeholder.text = str(raw_content)

        # ΑΥΤΟΜΑΤΗ ΕΙΚΟΝΑ: Αν το layout έχει Picture Placeholder (συνήθως ID 10-15 στο template σου)
        for placeholder in slide.placeholders:
            if (
                "Picture" in placeholder.name
                or placeholder.placeholder_format.type == 18
            ):  # 18 είναι το type για Picture
                img_url = get_unsplash_image(topic + " " + data.get("title", ""))
                if img_url:
                    img_path = download_image(img_url)
                    if img_path:
                        placeholder.insert_picture(img_path)
                        os.remove(img_path)  # Διαγραφή προσωρινού αρχείου

    # 4. Αποθήκευση: Κόβουμε το όνομα στους 30 χαρακτήρες και βγάζουμε τα σύμβολα
    safe_topic = (
        topic[:30].replace(" ", "_").replace(":", "").replace(",", "").replace(".", "")
    )
    output_name = f"Presentation_{safe_topic}.pptx"

    prs.save(output_name)
    print(f"\n--- Η ΠΑΡΟΥΣΙΑΣΗ ΕΙΝΑΙ ΕΤΟΙΜΗ: {output_name} ---")
    return output_name


if __name__ == "__main__":
    user_topic = input("Δώσε το θέμα της παρουσίασης: ")
    build_presentation(user_topic)
