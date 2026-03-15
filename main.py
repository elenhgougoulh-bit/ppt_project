from pptx import Presentation
import os
from dotenv import load_dotenv

# Φορτώνουμε τα κρυφά κλειδιά από το .env αρχείο
load_dotenv()

def create_test_presentation():
    print("Ξεκινάει η δημιουργία της παρουσίασης...")
    
    # 1. Δημιουργία μιας νέας (κενής) παρουσίασης
    prs = Presentation()
    
    # 2. Επιλογή του layout (0 = Κλασική διαφάνεια τίτλου)
    title_slide_layout = prs.slide_layouts[0]
    
    # 3. Δημιουργία της διαφάνειας
    slide = prs.slides.add_slide(title_slide_layout)
    
    # 4. Προσθήκη κειμένου στα "κουτιά" (placeholders)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Η Πρώτη μας AI Παρουσίαση!"
    subtitle.text = "Το Docker περιβάλλον δουλεύει άψογα."
    
    # 5. Αποθήκευση του αρχείου
    output_filename = 'test_output.pptx'
    prs.save(output_filename)
    
    print(f"Επιτυχία! Το αρχείο '{output_filename}' δημιουργήθηκε.")

if __name__ == "__main__":
    create_test_presentation()